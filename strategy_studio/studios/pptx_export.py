"""
PPTX Export Module — Board-ready PowerPoint generator.

Generates a consulting-grade slide deck from a StrategyReport analysis JSON.
Dark background, clean tables, minimal decoration.

Usage:
    from strategy_studio.studios.pptx_export import StrategyDeckGenerator
    gen = StrategyDeckGenerator()
    gen.generate(analysis_json, Path("output.pptx"))
"""

from __future__ import annotations

import datetime as dt
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Brand Colors ──────────────────────────────────────────────────────────
BG_COLOR = RGBColor(0x0A, 0x09, 0x07)
TEXT_COLOR = RGBColor(0xF2, 0xED, 0xE3)
ACCENT_COLOR = RGBColor(0xC8, 0xA9, 0x6E)
SIGNAL_RED = RGBColor(0xE0, 0x4F, 0x38)
SIGNAL_GREEN = RGBColor(0x4F, 0xB6, 0x8C)
SIGNAL_AMBER = RGBColor(0xF3, 0x9C, 0x12)
BORDER_COLOR = RGBColor(0x2A, 0x28, 0x25)
MUTED_TEXT = RGBColor(0x8A, 0x80, 0x75)
SECONDARY_BG = RGBColor(0x1A, 0x18, 0x15)
TIER_A = RGBColor(0x4F, 0xB6, 0x8C)
TIER_B = RGBColor(0x5D, 0x9B, 0xD4)
TIER_C = RGBColor(0xF3, 0x9C, 0x12)
TIER_D = RGBColor(0xE0, 0x4F, 0x38)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

_RECTANGLE = MSO_SHAPE_TYPE.RECTANGLE
_ROUNDED = MSO_SHAPE_TYPE.ROUNDED_RECTANGLE
_OVAL = MSO_SHAPE_TYPE.OVAL


def _confidence_color(conf: str) -> RGBColor:
    return {"H": SIGNAL_GREEN, "M": SIGNAL_AMBER, "L": SIGNAL_RED}.get(conf, SIGNAL_AMBER)


def _risk_color(level: str) -> RGBColor:
    return {
        "low": SIGNAL_GREEN,
        "medium": SIGNAL_AMBER,
        "high": SIGNAL_RED,
        "critical": RGBColor(0xC0, 0x00, 0x00),
    }.get(level, SIGNAL_AMBER)


def _tier_color(tier: str) -> RGBColor:
    return {"A": TIER_A, "B": TIER_B, "C": TIER_C, "D": TIER_D}.get(tier, TIER_C)


def _safe(obj, key: str, default: Any = "") -> Any:
    """Safely get from dict or return default."""
    if isinstance(obj, dict):
        return obj.get(key, default)
    return getattr(obj, key, default)


class StrategyDeckGenerator:
    """Generates board-ready PowerPoint from analysis JSON."""

    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H

    # ── Public API ──────────────────────────────────────────────────────

    def generate(self, analysis_json: dict, output_path: Path) -> Path:
        """Generate the full slide deck and save to output_path."""
        self._build_title_slide(analysis_json)
        self._build_executive_summary(analysis_json)
        self._build_decision_matrix(analysis_json)
        self._build_competitive_analysis(analysis_json)
        self._build_scenario_analysis(analysis_json)
        self._build_predictions(analysis_json)
        self._build_risks(analysis_json)
        self._build_next_steps(analysis_json)
        self.prs.save(str(output_path))
        return output_path

    # ── Slide helpers ───────────────────────────────────────────────────

    def _add_bg(self, slide) -> None:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def _add_footer(self, slide, page_num: int) -> None:
        txBox = slide.shapes.add_textbox(Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.35))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Strategy Studio  |  Confidential  |  Page {page_num}"
        p.font.size = Pt(9)
        p.font.color.rgb = MUTED_TEXT
        p.alignment = PP_ALIGN.RIGHT

    def _add_section_bar(self, slide) -> None:
        shape = slide.shapes.add_shape(
            _RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(0.06),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_COLOR
        shape.line.fill.background()

    def _slide_title(self, slide, title: str, top: float = 0.35) -> None:
        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1), Inches(0.7))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = TEXT_COLOR

    def _slide_subtitle(self, slide, text: str, top: float = 1.05) -> None:
        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1), Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(13)
        p.font.color.rgb = MUTED_TEXT

    def _blank_slide(self):
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        self._add_bg(slide)
        self._add_section_bar(slide)
        return slide

    # ── 1. Title Slide ──────────────────────────────────────────────────

    def _build_title_slide(self, data: dict) -> None:
        slide = self._blank_slide()
        es = data.get("executive_summary", {})

        company = data.get("company_name", data.get("title", "Strategy Analysis"))
        date_str = es.get("date", "")
        if not date_str:
            date_str = dt.datetime.now().strftime("%B %d, %Y")
        conf = es.get("confidence", "M")
        rec = es.get("recommendation", "")

        # Company name
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = company
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = TEXT_COLOR

        p2 = tf.add_paragraph()
        p2.text = "Strategy Analysis & Recommendations"
        p2.font.size = Pt(20)
        p2.font.color.rgb = ACCENT_COLOR
        p2.space_before = Pt(12)

        p3 = tf.add_paragraph()
        p3.text = date_str
        p3.font.size = Pt(14)
        p3.font.color.rgb = MUTED_TEXT
        p3.space_before = Pt(20)

        # Confidence badge
        badge_color = _confidence_color(conf)
        badge = slide.shapes.add_shape(
            _ROUNDED,
            Inches(4.5), Inches(4.0), Inches(2.2), Inches(0.55),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = badge_color
        badge.line.fill.background()
        tf_b = badge.text_frame
        tf_b.word_wrap = False
        tf_b.paragraphs[0].text = f"Confidence: {conf}"
        tf_b.paragraphs[0].font.size = Pt(14)
        tf_b.paragraphs[0].font.bold = True
        tf_b.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tf_b.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf_b.paragraphs[0].space_before = Pt(6)

        # Recommendation preview
        if rec:
            rec_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.0))
            rtf = rec_box.text_frame
            rtf.word_wrap = True
            rp = rtf.paragraphs[0]
            rp.text = rec[:200]
            rp.font.size = Pt(14)
            rp.font.color.rgb = TEXT_COLOR
            rp.font.italic = True

        self._add_footer(slide, 1)

    # ── 2. Executive Summary ────────────────────────────────────────────

    def _build_executive_summary(self, data: dict) -> None:
        slide = self._blank_slide()
        self._slide_title(slide, "Executive Summary")
        es = data.get("executive_summary", {})

        findings = es.get("key_findings", [])
        rec = es.get("recommendation", "")
        conf = es.get("confidence", "M")
        risks = es.get("risks", [])

        # Left column: Key Findings
        left = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.5))
        ltf = left.text_frame
        ltf.word_wrap = True

        hp = ltf.paragraphs[0]
        hp.text = "Key Findings"
        hp.font.size = Pt(16)
        hp.font.bold = True
        hp.font.color.rgb = ACCENT_COLOR
        hp.space_after = Pt(10)

        for finding in findings[:6]:
            p = ltf.add_paragraph()
            p.text = finding[:150]
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_COLOR
            p.space_before = Pt(6)
            p.level = 0

        # Right column: Recommendation + Confidence
        right = slide.shapes.add_textbox(Inches(6.8), Inches(1.3), Inches(5.9), Inches(5.5))
        rtf = right.text_frame
        rtf.word_wrap = True

        rp = rtf.paragraphs[0]
        rp.text = "Primary Recommendation"
        rp.font.size = Pt(16)
        rp.font.bold = True
        rp.font.color.rgb = ACCENT_COLOR
        rp.space_after = Pt(8)

        rp2 = rtf.add_paragraph()
        rp2.text = rec[:300] if rec else "No recommendation provided."
        rp2.font.size = Pt(13)
        rp2.font.color.rgb = TEXT_COLOR
        rp2.space_after = Pt(16)

        rp3 = rtf.add_paragraph()
        rp3.text = f"Confidence Level: {conf}"
        rp3.font.size = Pt(14)
        rp3.font.bold = True
        rp3.font.color.rgb = _confidence_color(conf)
        rp3.space_after = Pt(16)

        if risks:
            rp4 = rtf.add_paragraph()
            rp4.text = "Key Risks"
            rp4.font.size = Pt(14)
            rp4.font.bold = True
            rp4.font.color.rgb = SIGNAL_RED
            rp4.space_after = Pt(6)
            for risk in risks[:4]:
                rp5 = rtf.add_paragraph()
                rp5.text = risk[:120]
                rp5.font.size = Pt(11)
                rp5.font.color.rgb = MUTED_TEXT
                rp5.space_before = Pt(4)

        self._add_footer(slide, 2)

    # ── 3. Decision Matrix ──────────────────────────────────────────────

    def _build_decision_matrix(self, data: dict) -> None:
        slide = self._blank_slide()
        self._slide_title(slide, "Decision Matrix")

        dr = data.get("decision_room", {})
        dm = dr.get("decision_matrix", {}) if isinstance(dr, dict) else {}
        options = dm.get("options", []) if isinstance(dm, dict) else []

        if not options:
            self._slide_subtitle(slide, "No decision matrix data available.")
            self._add_footer(slide, 3)
            return

        rows = min(len(options) + 1, 9)
        cols_count = 5
        left = Inches(0.6)
        top = Inches(1.3)
        width = Inches(8.0)
        height = Inches(5.5)

        table = slide.shapes.add_table(rows, cols_count, left, top, width, height).table
        headers = ["Rank", "Option", "Score", "Tier", "Conf."]
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = SECONDARY_BG
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.font.bold = True
            para.font.color.rgb = ACCENT_COLOR

        for row_idx, opt in enumerate(options[: rows - 1]):
            r = row_idx + 1
            vals = [
                f"#{opt.get('rank', r)}",
                opt.get("option_title", opt.get("title", f"Option {r}"))[:40],
                f"{opt.get('total_score', 0):.2f}",
                opt.get("tier", "C"),
                opt.get("confidence", "M"),
            ]
            for col_idx, val in enumerate(vals):
                cell = table.cell(r, col_idx)
                cell.text = val
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(11)
                para.font.color.rgb = TEXT_COLOR
                if col_idx == 3:
                    para.font.color.rgb = _tier_color(val)
                    para.font.bold = True
                if col_idx == 0:
                    para.font.color.rgb = ACCENT_COLOR

        # Bar chart on right
        chart_left = Inches(9.0)
        chart_top = Inches(1.3)
        chart_width = Inches(3.8)
        chart_height = Inches(4.0)

        chart_data = ChartData()
        chart_data.categories = [
            opt.get("option_title", f"Opt{i+1}")[:20]
            for i, opt in enumerate(options[:6])
        ]
        scores = [opt.get("total_score", 0) for opt in options[:6]]
        chart_data.add_series("Score", scores)

        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            chart_left, chart_top, chart_width, chart_height,
            chart_data,
        )
        graphic_frame.chart.has_legend = False

        self._add_footer(slide, 3)

    # ── 4. Competitive Analysis ─────────────────────────────────────────

    def _build_competitive_analysis(self, data: dict) -> None:
        slide = self._blank_slide()
        self._slide_title(slide, "Competitive Analysis")

        wg = data.get("wargame_result", {})
        if not wg:
            self._slide_subtitle(slide, "No competitive analysis data available.")
            self._add_footer(slide, 4)
            return

        scenario = wg.get("scenario_name", "Competitive Landscape")
        risk = wg.get("risk_level", "medium")
        moves = wg.get("moves", [])
        response = wg.get("recommended_response", "")

        # Scenario header
        hdr = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(8.0), Inches(0.5))
        htf = hdr.text_frame
        hp = htf.paragraphs[0]
        hp.text = scenario
        hp.font.size = Pt(18)
        hp.font.bold = True
        hp.font.color.rgb = TEXT_COLOR

        # Risk badge
        rc = _risk_color(risk)
        badge = slide.shapes.add_shape(
            _ROUNDED,
            Inches(9.0), Inches(1.2), Inches(1.8), Inches(0.45),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = rc
        badge.line.fill.background()
        btf = badge.text_frame
        btf.word_wrap = False
        bp = btf.paragraphs[0]
        bp.text = risk.upper()
        bp.font.size = Pt(12)
        bp.font.bold = True
        bp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        bp.alignment = PP_ALIGN.CENTER

        # Moves table
        move_rows = min(len(moves) + 1, 8)
        mtable = slide.shapes.add_table(
            move_rows, 4,
            Inches(0.6), Inches(1.9), Inches(8.0), Inches(4.5),
        ).table
        for i, h in enumerate(["Actor", "Move", "Prob.", "Impact"]):
            cell = mtable.cell(0, i)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = SECONDARY_BG
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR

        for ri, move in enumerate(moves[: move_rows - 1]):
            r = ri + 1
            actor = move.get("actor", "")
            move_text = move.get("move", "")[:60]
            prob = move.get("probability", 0)
            impact = move.get("impact", "")[:40]
            for ci, val in enumerate([actor, move_text, f"{prob:.0%}", impact]):
                cell = mtable.cell(r, ci)
                cell.text = val
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR

        # Recommended response
        if response:
            resp_box = slide.shapes.add_textbox(Inches(0.6), Inches(5.8), Inches(12.1), Inches(0.8))
            rtf = resp_box.text_frame
            rtf.word_wrap = True
            rp = rtf.paragraphs[0]
            rp.text = "Recommended Response"
            rp.font.size = Pt(13)
            rp.font.bold = True
            rp.font.color.rgb = SIGNAL_GREEN
            rp2 = rtf.add_paragraph()
            rp2.text = response[:250]
            rp2.font.size = Pt(11)
            rp2.font.color.rgb = TEXT_COLOR

        self._add_footer(slide, 4)

    # ── 5. Scenario Analysis ────────────────────────────────────────────

    def _build_scenario_analysis(self, data: dict) -> None:
        slide = self._blank_slide()
        self._slide_title(slide, "Scenario Analysis")

        dr = data.get("decision_room", {})
        scenarios = dr.get("scenarios_considered", []) if isinstance(dr, dict) else []

        if not scenarios:
            self._slide_subtitle(slide, "No scenario data available.")
            self._add_footer(slide, 5)
            return

        cols = 3
        card_w = Inches(3.8)
        card_h = Inches(2.8)
        gap_x = Inches(0.3)
        gap_y = Inches(0.2)
        start_left = Inches(0.6)
        start_top = Inches(1.3)

        for idx, sc in enumerate(scenarios[:6]):
            col = idx % cols
            row = idx // cols
            left = start_left + col * (card_w + gap_x)
            top = start_top + row * (card_h + gap_y)

            # Card background
            card = slide.shapes.add_shape(_RECTANGLE, left, top, card_w, card_h)
            card.fill.solid()
            card.fill.fore_color.rgb = SECONDARY_BG
            card.line.color.rgb = BORDER_COLOR
            card.line.width = Pt(0.5)

            # Card content
            txBox = slide.shapes.add_textbox(
                left + Inches(0.15), top + Inches(0.1),
                card_w - Inches(0.3), card_h - Inches(0.2),
            )
            tf = txBox.text_frame
            tf.word_wrap = True

            name = sc.get("name", f"Scenario {idx+1}")
            prob = sc.get("probability", 0)
            assumptions = sc.get("assumptions", [])
            variables = sc.get("variables", {})

            p = tf.paragraphs[0]
            p.text = name[:35]
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = ACCENT_COLOR
            p.space_after = Pt(4)

            bar_filled = int(prob * 20)
            bar = "\u2588" * bar_filled + "\u2591" * (20 - bar_filled)
            p2 = tf.add_paragraph()
            p2.text = f"{bar}  {prob:.0%}"
            p2.font.size = Pt(10)
            p2.font.color.rgb = SIGNAL_GREEN
            p2.space_after = Pt(6)

            if assumptions:
                p3 = tf.add_paragraph()
                p3.text = "Assumptions:"
                p3.font.size = Pt(9)
                p3.font.bold = True
                p3.font.color.rgb = MUTED_TEXT
                for assump in assumptions[:2]:
                    p4 = tf.add_paragraph()
                    p4.text = f"  {assump[:50]}"
                    p4.font.size = Pt(9)
                    p4.font.color.rgb = MUTED_TEXT

            if variables:
                p5 = tf.add_paragraph()
                p5.text = "Variables:"
                p5.font.size = Pt(9)
                p5.font.bold = True
                p5.font.color.rgb = MUTED_TEXT
                for k, v in list(variables.items())[:2]:
                    p6 = tf.add_paragraph()
                    p6.text = f"  {k}: {v:.2f}"
                    p6.font.size = Pt(9)
                    p6.font.color.rgb = MUTED_TEXT

        self._add_footer(slide, 5)

    # ── 6. Predictions ──────────────────────────────────────────────────

    def _build_predictions(self, data: dict) -> None:
        slide = self._blank_slide()
        self._slide_title(slide, "Predictions")

        preds = []
        pr = data.get("prediction_result", {})
        if pr:
            if isinstance(pr, list):
                preds = pr
            elif isinstance(pr, dict) and pr:
                preds = [pr]

        if not preds:
            self._slide_subtitle(slide, "No prediction data available.")
            self._add_footer(slide, 6)
            return

        rows = min(len(preds) + 1, 9)
        ptable = slide.shapes.add_table(
            rows, 5,
            Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.5),
        ).table
        for i, h in enumerate(["Variable", "Point Est.", "CI Low", "CI High", "Method"]):
            cell = ptable.cell(0, i)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = SECONDARY_BG
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR

        for ri, pred in enumerate(preds[: rows - 1]):
            r = ri + 1
            ci = pred.get("confidence_interval", [0, 0])
            if isinstance(ci, (list, tuple)) and len(ci) == 2:
                ci_low, ci_high = ci[0], ci[1]
            else:
                ci_low, ci_high = 0, 0
            vals = [
                pred.get("variable", f"Var {r}")[:25],
                f"{pred.get('point_estimate', 0):,.2f}",
                f"{ci_low:,.2f}",
                f"{ci_high:,.2f}",
                pred.get("method", "")[:20],
            ]
            for ci2, val in enumerate(vals):
                cell = ptable.cell(r, ci2)
                cell.text = val
                cell.text_frame.paragraphs[0].font.size = Pt(11)
                cell.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR

        self._add_footer(slide, 6)

    # ── 7. Risks ────────────────────────────────────────────────────────

    def _build_risks(self, data: dict) -> None:
        slide = self._blank_slide()
        self._slide_title(slide, "Risk Assessment")

        es = data.get("executive_summary", {})
        risks = es.get("risks", [])
        dr = data.get("decision_room", {})
        dr_risks = dr.get("risks", []) if isinstance(dr, dict) else []
        all_risks = risks + dr_risks

        if not all_risks:
            self._slide_subtitle(slide, "No risk data available.")
            self._add_footer(slide, 7)
            return

        rows = min(len(all_risks) + 1, 10)
        rtable = slide.shapes.add_table(
            rows, 4,
            Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.5),
        ).table
        for i, h in enumerate(["#", "Risk", "Severity", "Mitigation"]):
            cell = rtable.cell(0, i)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = SECONDARY_BG
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR

        for ri, risk in enumerate(all_risks[: rows - 1]):
            r = ri + 1
            if isinstance(risk, dict):
                risk_text = risk.get("description", risk.get("risk", str(risk)))[:80]
                severity = risk.get("severity", "medium")
                mitigation = risk.get("mitigation", "Monitor and reassess")[:60]
            else:
                risk_text = str(risk)[:80]
                severity = "medium"
                mitigation = "Monitor and reassess"

            for ci, val in enumerate([str(r), risk_text, severity.capitalize(), mitigation]):
                cell = rtable.cell(r, ci)
                cell.text = val
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
                if ci == 2:
                    cell.text_frame.paragraphs[0].font.color.rgb = _risk_color(severity)
                    cell.text_frame.paragraphs[0].font.bold = True

        self._add_footer(slide, 7)

    # ── 8. Next Steps ───────────────────────────────────────────────────

    def _build_next_steps(self, data: dict) -> None:
        slide = self._blank_slide()
        self._slide_title(slide, "Next Steps")

        es = data.get("executive_summary", {})
        steps = es.get("next_steps", [])
        dr = data.get("decision_room", {})
        dr_steps = dr.get("next_steps", []) if isinstance(dr, dict) else []
        all_steps = steps + dr_steps

        if not all_steps:
            self._slide_subtitle(slide, "No next steps defined.")
            self._add_footer(slide, 8)
            return

        top = Inches(1.3)
        for i, step in enumerate(all_steps[:8]):
            step_text = str(step)[:150] if not isinstance(step, str) else step[:150]

            # Number circle
            circle = slide.shapes.add_shape(
                _OVAL,
                Inches(0.8), top + Inches(i * 0.65), Inches(0.4), Inches(0.4),
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = ACCENT_COLOR
            circle.line.fill.background()
            ctf = circle.text_frame
            ctf.word_wrap = False
            cp = ctf.paragraphs[0]
            cp.text = str(i + 1)
            cp.font.size = Pt(12)
            cp.font.bold = True
            cp.font.color.rgb = BG_COLOR
            cp.alignment = PP_ALIGN.CENTER
            cp.space_before = Pt(2)

            # Step text
            txBox = slide.shapes.add_textbox(
                Inches(1.4), top + Inches(i * 0.65) + Inches(0.05),
                Inches(11.0), Inches(0.5),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = step_text
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_COLOR

        self._add_footer(slide, 8)
